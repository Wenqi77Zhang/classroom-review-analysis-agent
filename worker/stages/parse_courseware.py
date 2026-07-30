"""Extract page-aligned text from PDF and PPTX courseware."""

from __future__ import annotations

from pathlib import Path
from uuid import UUID

from pptx import Presentation
from pypdf import PdfReader

from worker.courseware_types import CoursewareDocument, CoursewarePage
from worker.errors import WorkerError, WorkerErrorCode

PDF_CONTENT_TYPE = "application/pdf"
PPTX_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


def _normalize_text(parts: list[str]) -> str:
    lines: list[str] = []
    for part in parts:
        for raw_line in part.splitlines():
            line = " ".join(raw_line.split())
            if line:
                lines.append(line)
    return "\n".join(lines)


def _parse_pdf(path: Path, *, asset_id: UUID, max_pages: int) -> CoursewareDocument:
    reader = PdfReader(path)
    if reader.is_encrypted:
        raise ValueError("encrypted PDF")
    if not reader.pages or len(reader.pages) > max_pages:
        raise ValueError("invalid PDF page count")
    pages = tuple(
        CoursewarePage(
            page_no=index,
            text=_normalize_text([page.extract_text() or ""]),
        )
        for index, page in enumerate(reader.pages, start=1)
    )
    return CoursewareDocument(asset_id=asset_id, pages=pages)


def _parse_pptx(path: Path, *, asset_id: UUID, max_pages: int) -> CoursewareDocument:
    presentation = Presentation(path)
    if not presentation.slides or len(presentation.slides) > max_pages:
        raise ValueError("invalid PPTX slide count")

    pages: list[CoursewarePage] = []
    for page_no, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    parts.extend(cell.text for cell in row.cells)
            elif shape.has_text_frame:
                parts.extend(paragraph.text for paragraph in shape.text_frame.paragraphs)
        if slide.has_notes_slide:
            parts.append(slide.notes_slide.notes_text_frame.text)
        pages.append(CoursewarePage(page_no=page_no, text=_normalize_text(parts)))
    return CoursewareDocument(asset_id=asset_id, pages=tuple(pages))


def parse_courseware(
    path: Path,
    *,
    asset_id: UUID,
    content_type: str,
    max_pages: int = 500,
) -> CoursewareDocument:
    """Dispatch on trusted MIME and return stable public failures."""

    if content_type not in {PDF_CONTENT_TYPE, PPTX_CONTENT_TYPE}:
        raise WorkerError(
            WorkerErrorCode.COURSEWARE_UNSUPPORTED,
            "课件类型不在 Worker 支持范围内。",
            retryable=False,
        )
    if max_pages < 1:
        raise WorkerError(
            WorkerErrorCode.COURSEWARE_PARSE_FAILED,
            "课件页数限制无效。",
            retryable=False,
        )

    try:
        if content_type == PDF_CONTENT_TYPE:
            return _parse_pdf(path, asset_id=asset_id, max_pages=max_pages)
        return _parse_pptx(path, asset_id=asset_id, max_pages=max_pages)
    except WorkerError:
        raise
    except Exception:  # noqa: BLE001 - parser libraries expose varied private exceptions
        raise WorkerError(
            WorkerErrorCode.COURSEWARE_PARSE_FAILED,
            "课件无法解析为可引用的页级文本。",
            retryable=False,
        ) from None
