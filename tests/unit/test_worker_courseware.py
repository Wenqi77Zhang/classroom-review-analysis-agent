from __future__ import annotations

from pathlib import Path
from uuid import uuid4
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter
from reportlab.pdfgen import canvas

from worker.courseware_types import CoursewareDocument, CoursewarePage
from worker.errors import WorkerError, WorkerErrorCode
from worker.stages.parse_courseware import (
    PDF_CONTENT_TYPE,
    PPTX_CONTENT_TYPE,
    parse_courseware,
)


def _write_pdf(path: Path) -> None:
    document = canvas.Canvas(str(path))
    document.drawString(72, 720, "Artificial intelligence")
    document.showPage()
    document.drawString(72, 720, "Machine learning")
    document.save()


def _write_pptx(path: Path) -> None:
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[5])
    first.shapes.title.text = "Algorithm"
    table = first.shapes.add_table(
        1,
        1,
        Inches(1),
        Inches(2),
        Inches(4),
        Inches(1),
    ).table
    table.cell(0, 0).text = "Complexity"
    first.notes_slide.notes_text_frame.text = "Instructor note"

    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "Humanities"
    presentation.save(path)


def test_courseware_types_use_one_based_pages() -> None:
    with pytest.raises(ValueError):
        CoursewarePage(page_no=0, text="invalid")

    page = CoursewarePage(page_no=1, text="Introduction")
    document = CoursewareDocument(asset_id=uuid4(), pages=(page,))

    assert document.pages[0].page_no == 1


def test_parse_pdf_extracts_text_and_one_based_pages(tmp_path: Path) -> None:
    path = tmp_path / "source-without-trusted-suffix"
    _write_pdf(path)
    asset_id = uuid4()

    document = parse_courseware(
        path,
        asset_id=asset_id,
        content_type=PDF_CONTENT_TYPE,
    )

    assert document.asset_id == asset_id
    assert [page.page_no for page in document.pages] == [1, 2]
    assert "Artificial intelligence" in document.pages[0].text
    assert "Machine learning" in document.pages[1].text


def test_parse_pptx_extracts_visible_shapes_and_tables_but_excludes_notes(
    tmp_path: Path,
) -> None:
    path = tmp_path / "source-without-trusted-suffix"
    _write_pptx(path)

    document = parse_courseware(
        path,
        asset_id=uuid4(),
        content_type=PPTX_CONTENT_TYPE,
    )

    assert [page.page_no for page in document.pages] == [1, 2]
    assert document.pages[0].text.splitlines() == [
        "Algorithm",
        "Complexity",
    ]
    assert "Instructor note" not in document.pages[0].text
    assert document.pages[1].text == "Humanities"


def test_parse_pptx_does_not_create_missing_notes_slide(tmp_path: Path) -> None:
    source = tmp_path / "source.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "No notes"
    assert slide.has_notes_slide is False
    presentation.save(source)

    parse_courseware(
        source,
        asset_id=uuid4(),
        content_type=PPTX_CONTENT_TYPE,
    )

    reopened = Presentation(source)
    assert reopened.slides[0].has_notes_slide is False


def test_parse_courseware_rejects_unsupported_mime(tmp_path: Path) -> None:
    path = tmp_path / "courseware.txt"
    path.write_text("content", encoding="utf-8")

    with pytest.raises(WorkerError) as raised:
        parse_courseware(
            path,
            asset_id=uuid4(),
            content_type="text/plain",
        )

    assert raised.value.code is WorkerErrorCode.COURSEWARE_UNSUPPORTED
    assert raised.value.retryable is False


def test_parse_courseware_rejects_broken_pdf_without_path_leak(tmp_path: Path) -> None:
    path = tmp_path / "private-courseware.pdf"
    path.write_bytes(b"not a pdf")

    with pytest.raises(WorkerError) as raised:
        parse_courseware(
            path,
            asset_id=uuid4(),
            content_type=PDF_CONTENT_TYPE,
        )

    assert raised.value.code is WorkerErrorCode.COURSEWARE_PARSE_FAILED
    assert raised.value.retryable is False
    assert str(path) not in str(raised.value)


def test_parse_courseware_rejects_encrypted_pdf(tmp_path: Path) -> None:
    path = tmp_path / "encrypted"
    writer = PdfWriter()
    writer.add_blank_page(width=300, height=300)
    writer.encrypt("secret")
    with path.open("wb") as output:
        writer.write(output)

    with pytest.raises(WorkerError) as raised:
        parse_courseware(
            path,
            asset_id=uuid4(),
            content_type=PDF_CONTENT_TYPE,
        )

    assert raised.value.code is WorkerErrorCode.COURSEWARE_PARSE_FAILED


def test_parse_courseware_rejects_broken_pptx(tmp_path: Path) -> None:
    path = tmp_path / "broken"
    path.write_bytes(b"not a pptx")

    with pytest.raises(WorkerError) as raised:
        parse_courseware(
            path,
            asset_id=uuid4(),
            content_type=PPTX_CONTENT_TYPE,
        )

    assert raised.value.code is WorkerErrorCode.COURSEWARE_PARSE_FAILED


def test_parse_courseware_rejects_excessive_pptx_compression_ratio(
    tmp_path: Path,
) -> None:
    path = tmp_path / "compressed-bomb"
    with ZipFile(path, "w", compression=ZIP_DEFLATED) as archive:
        archive.writestr("oversized.bin", b"0" * (2 * 1024 * 1024))

    with pytest.raises(WorkerError) as raised:
        parse_courseware(
            path,
            asset_id=uuid4(),
            content_type=PPTX_CONTENT_TYPE,
        )

    assert raised.value.code is WorkerErrorCode.COURSEWARE_PARSE_FAILED


@pytest.mark.parametrize("content_type", [PDF_CONTENT_TYPE, PPTX_CONTENT_TYPE])
def test_parse_courseware_enforces_page_limit(
    tmp_path: Path,
    content_type: str,
) -> None:
    path = tmp_path / "too-many-pages"
    if content_type == PDF_CONTENT_TYPE:
        _write_pdf(path)
    else:
        _write_pptx(path)

    with pytest.raises(WorkerError) as raised:
        parse_courseware(
            path,
            asset_id=uuid4(),
            content_type=content_type,
            max_pages=1,
        )

    assert raised.value.code is WorkerErrorCode.COURSEWARE_PARSE_FAILED


def test_parse_courseware_rejects_zero_page_pdf(tmp_path: Path) -> None:
    path = tmp_path / "empty"
    writer = PdfWriter()
    with path.open("wb") as output:
        writer.write(output)

    with pytest.raises(WorkerError) as raised:
        parse_courseware(
            path,
            asset_id=uuid4(),
            content_type=PDF_CONTENT_TYPE,
        )

    assert raised.value.code is WorkerErrorCode.COURSEWARE_PARSE_FAILED
